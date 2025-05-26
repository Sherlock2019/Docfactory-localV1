import re
import os
from io import BytesIO
from datetime import date
from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches as PptxInches # For PPTX image sizing
import pandas as pd
import uuid # For temporary filenames

from utils.extract_placeholders import extract_placeholders, extract_text_from_file

def fill_document_template(
    tpl_path: str,
    file_inputs: dict[str, "werkzeug.datastructures.FileStorage"], # Pass FileStorage object directly
    text_inputs: dict[str, str],
    out_path: str
):
    """
    Fills a .docx or .pptx template using both text and file-based placeholder values.

    Args:
    - tpl_path:      Path to the input .docx or .pptx template
    - file_inputs:   Dict mapping placeholder -> FileStorage (uploaded file)
    - text_inputs:   Dict mapping placeholder -> raw string (manual entry)
    - out_path:      Where to write the final filled document
    """
    file_extension = tpl_path.lower().rsplit(".", 1)[-1]
    today_str = date.today().strftime("%Y-%m-%d") # Use hyphens for readability

    if file_extension == "docx":
        doc = Document(tpl_path)
        placeholders = extract_placeholders(tpl_path)

        for ph in placeholders:
            tag = f"{{{ph}}}"
            file_storage = file_inputs.get(ph) # This is now the FileStorage object
            value = text_inputs.get(ph, "")

            # Iterate over all paragraphs and runs to find and replace
            for para in doc.paragraphs:
                # Optimized replacement for text placeholders (in-place)
                if tag in para.text:
                    # Clean up the paragraph's text before inserting new runs if needed
                    # This advanced replacement ensures surrounding text is kept.
                    # This part is simplified for brevity; a full solution is more robust.
                    # For complex scenarios like multiple placeholders in one run,
                    # or partial replacement within a run, more logic is needed.
                    inline_replacement_done = False
                    if not file_storage: # Only for pure text replacement
                        # Replace the tag directly in the paragraph's text
                        # This works if the entire paragraph is the placeholder or if
                        # you don't care about run-level formatting.
                        # For true in-place formatting, you'd iterate runs and replace there.
                        para.text = para.text.replace(tag, str(value))
                        inline_replacement_done = True

                    if not inline_replacement_done: # If we need to insert objects or large text
                        # Remove the placeholder text from the paragraph
                        para.text = para.text.replace(tag, "")

                        # Handle IMAGE upload
                        if file_storage and file_storage.filename.lower().endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp")):
                            img_data = file_storage.read()
                            tmp_path = f"/tmp/{uuid.uuid4()}_{file_storage.filename}" # Unique temporary name
                            with open(tmp_path, "wb") as tmp:
                                tmp.write(img_data)
                            para.add_run().add_picture(tmp_path, width=Inches(4))
                            os.remove(tmp_path)
                            file_storage.stream.seek(0) # Reset stream if read for image processing

                        # Handle Excel (.xlsx) as table
                        elif file_storage and file_storage.filename.lower().endswith(".xlsx"):
                            file_storage.stream.seek(0) # Ensure stream is at the beginning
                            try:
                                df = pd.read_excel(file_storage.stream)
                                # Add table at the end of the document (python-docx limitation for in-place)
                                # To insert truly in-place, you'd need to delete the placeholder paragraph
                                # and insert the table before/after a specific content control or bookmark.
                                table = doc.add_table(rows=1, cols=len(df.columns))
                                hdr_cells = table.rows[0].cells
                                for i, col in enumerate(df.columns):
                                    hdr_cells[i].text = str(col)
                                for _, row in df.iterrows():
                                    row_cells = table.add_row().cells
                                    for i, cell in enumerate(row):
                                        row_cells[i].text = str(cell)
                            except Exception as e:
                                para.add_run(f"[Error reading Excel: {e}]")
                            file_storage.stream.seek(0) # Reset stream after reading

                        # Handle txt/docx/pptx as plain text insertion
                        elif file_storage:
                            file_storage.stream.seek(0) # Ensure stream is at the beginning
                            extracted = extract_text_from_file(file_storage.stream, file_storage.filename)
                            # Replace placeholder with extracted text, adding line breaks
                            # This replaces the original paragraph content if the placeholder was found.
                            # For in-place text: consider splitting the paragraph into runs and inserting.
                            # For now, it will replace the paragraph with the extracted text.
                            for line in extracted.splitlines():
                                para.add_run(line).add_break()
                            file_storage.stream.seek(0) # Reset stream after reading

                        # Manual text input fallback (if not handled by inline replacement above)
                        # This would only be hit if a placeholder was in the paragraph,
                        # but it wasn't a file input and wasn't handled by inline_replacement_done
                        else:
                            # This block is now less likely to be hit for pure text
                            # placeholders if the inline replacement logic works.
                            para.add_run(str(value))

            # Iterate over all tables in the document to find and replace in cells
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if tag in cell.text:
                            # Replace the tag directly in the cell's text
                            cell.text = cell.text.replace(tag, str(value))


        doc.save(out_path)
    elif file_extension == "pptx":
        prs = Presentation(tpl_path)
        placeholders = extract_placeholders(tpl_path)

        for ph in placeholders:
            tag = f"{{{ph}}}"
            file_storage = file_inputs.get(ph)
            value = text_inputs.get(ph, "")

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        if tag in shape.text_frame.text:
                            # PPTX: Simple text replacement is possible.
                            # In-place image/table insertion is significantly harder with python-pptx
                            # and often requires more complex slide manipulation or pre-defined placeholders.
                            # We'll prioritize text replacement.
                            if file_storage:
                                if file_storage.filename.lower().endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp")):
                                    # PPTX: You can add pictures, but not easily in place of a placeholder.
                                    # For a simple solution, we'll replace the text placeholder with a notice.
                                    # For real image insertion, you'd typically need to design a placeholder
                                    # picture shape in the template and replace its content.
                                    shape.text_frame.text = shape.text_frame.text.replace(tag, "[IMAGE PLACEHOLDER - INSERT MANUALLY]")
                                    # Alternatively, add a new picture to the slide at a default position:
                                    # slide.shapes.add_picture(file_storage.stream, PptxInches(1), PptxInches(1))
                                    # But this won't be "in-place".
                                else: # Handle other file types as text
                                    file_storage.stream.seek(0)
                                    extracted = extract_text_from_file(file_storage.stream, file_storage.filename)
                                    shape.text_frame.text = shape.text_frame.text.replace(tag, extracted)
                                file_storage.stream.seek(0) # Reset stream after reading
                            else:
                                shape.text_frame.text = shape.text_frame.text.replace(tag, str(value))
        prs.save(out_path)
    else:
        raise ValueError(f"Unsupported template file type: .{file_extension}")


