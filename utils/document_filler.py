import re
import os
from io import BytesIO
from datetime import date
from docx import Document
from docx.shared import Inches
from docx.oxml.ns import qn # Needed for XML manipulation
from docx.oxml import OxmlElement # Needed for XML manipulation
# No need for direct Paragraph or Table imports if iterating doc.paragraphs and doc.tables
from pptx import Presentation
from pptx.util import Inches as PptxInches # For PPTX image sizing
import pandas as pd
import uuid # For temporary filenames

from utils.extract_placeholders import extract_placeholders, extract_text_from_file

def fill_document_template(
    tpl_path: str,
    file_inputs: dict[str, "werkzeug.datastructures.FileStorage"],
    text_inputs: dict[str, str],
    out_path: str
):
    """
    Fills a .docx or .pptx template using both text and file-based placeholder values.

    Args:
    - tpl_path:      Path to the input .docx or .pptx template
    - file_inputs:   Dict mapping placeholder -> FileStorage (uploaded file)
    - text_inputs:   Dict mapping placeholder -> raw string (manual entry)
    - out_path:      Where to write the final filled document
    """
    file_extension = tpl_path.lower().rsplit(".", 1)[-1]
    today_str = date.today().strftime("%Y-%m-%d")

    if file_extension == "docx":
        doc = Document(tpl_path)
        
        # A list to keep track of paragraphs that need to be removed (e.g., replaced by tables)
        paragraphs_to_remove = []

        # First pass: Handle text placeholders and mark paragraphs for removal if replaced by a file
        for p in doc.paragraphs:
            for ph_name in list(file_inputs.keys()) + list(text_inputs.keys()):
                tag = f"{{{ph_name}}}"
                if tag in p.text:
                    file_storage = file_inputs.get(ph_name)
                    value = text_inputs.get(ph_name, "")

                    if not file_storage: # Pure text replacement
                        p.text = p.text.replace(tag, str(value))
                    else: # File-based content (Image, Excel, Text file)
                        # Mark this paragraph for removal *if* we are inserting a new element
                        # like a table that replaces it.
                        # For images/text files, we can often modify the paragraph directly.
                        if file_storage.filename.lower().endswith(".xlsx"):
                            paragraphs_to_remove.append(p)
                        
                        # Remove the placeholder text from the paragraph's current content
                        # This prepares the paragraph to be either replaced, or to have its content cleared for image/text
                        p.text = p.text.replace(tag, "")
                        
                        # --- Handle IMAGE upload ---
                        if file_storage.filename.lower().endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp")):
                            img_data = file_storage.read()
                            tmp_path = f"/tmp/{uuid.uuid4()}_{file_storage.filename}"
                            with open(tmp_path, "wb") as tmp:
                                tmp.write(img_data)
                            
                            # Clear current runs in paragraph and add image
                            # This replaces the content of the existing paragraph with the image
                            for run in p.runs:
                                run.clear()
                            p.add_run().add_picture(tmp_path, width=Inches(4))
                            os.remove(tmp_path)
                            file_storage.stream.seek(0)


                        # --- Handle Excel (.xlsx) as table ---
                        elif file_storage.filename.lower().endswith(".xlsx"):
                            file_storage.stream.seek(0)
                            try:
                                df = pd.read_excel(file_storage.stream)
                                
                                # Get the XML element of the paragraph we are replacing
                                p_xml_element = p._element
                                # Get its parent (the document body XML element)
                                body_xml_element = p_xml_element.getparent()
                                
                                # Create new table XML element directly using python-docx's methods
                                # This is the more robust way to create a table that can be inserted
                                new_table = doc.add_table(rows=1, cols=len(df.columns)) # Temporarily add to end
                                
                                # Now, move the newly created table's XML element to the correct position
                                # This is where the magic happens: get the XML of the new table
                                # and insert it before the original placeholder paragraph's XML.
                                new_table_xml = new_table._element
                                body_xml_element.insert(body_xml_element.index(p_xml_element), new_table_xml)

                                # Fill the table with data
                                hdr_cells = new_table.rows[0].cells
                                for i, col in enumerate(df.columns):
                                    hdr_cells[i].text = str(col)
                                for _, row in df.iterrows():
                                    row_cells = new_table.add_row().cells
                                    for i, cell in enumerate(row):
                                        row_cells[i].text = str(cell)
                                
                            except Exception as e:
                                # If error, replace placeholder paragraph content with error message
                                p.text = f"[Error reading Excel for {ph_name}: {e}]"
                                # No need to remove the paragraph if it now contains an error message.
                                # Remove it from paragraphs_to_remove if it was added.
                                if p in paragraphs_to_remove:
                                    paragraphs_to_remove.remove(p)

                            file_storage.stream.seek(0)


                        # --- Handle txt/docx/pptx as plain text insertion ---
                        elif file_storage:
                            file_storage.stream.seek(0)
                            extracted = extract_text_from_file(file_storage.stream, file_storage.filename)
                            
                            # Replace content of placeholder paragraph
                            for run in p.runs:
                                run.clear()
                            for line in extracted.splitlines():
                                p.add_run(line)
                                p.add_run().add_break()
                            file_storage.stream.seek(0)
                        
                        else: # Fallback for other manual text inputs if somehow not caught
                             p.text = p.text.replace(tag, str(value))
                    break # Break after handling a placeholder in this paragraph

        # Second pass: Handle text placeholders in tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for ph_name in text_inputs.keys():
                        tag = f"{{{ph_name}}}"
                        if tag in cell.text:
                            cell.text = cell.text.replace(tag, str(text_inputs[ph_name]))

        # Finally, remove the marked placeholder paragraphs
        for p_to_remove in paragraphs_to_remove:
            p_to_remove._element.getparent().remove(p_to_remove._element)

        doc.save(out_path)

    elif file_extension == "pptx":
        prs = Presentation(tpl_path)
        placeholders = extract_placeholders(tpl_path)

        for ph_name in placeholders:
            tag = f"{{{ph_name}}}"
            file_storage = file_inputs.get(ph_name)
            value = text_inputs.get(ph_name, "")

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        if tag in shape.text_frame.text:
                            if file_storage:
                                if file_storage.filename.lower().endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp")):
                                    shape.text_frame.text = shape.text_frame.text.replace(tag, "[IMAGE PLACEHOLDER - INSERT MANUALLY]")
                                else:
                                    file_storage.stream.seek(0)
                                    extracted = extract_text_from_file(file_storage.stream, file_storage.filename)
                                    shape.text_frame.text = shape.text_frame.text.replace(tag, extracted)
                                file_storage.stream.seek(0)
                            else:
                                shape.text_frame.text = shape.text_frame.text.replace(tag, str(value))
        prs.save(out_path)
    else:
        raise ValueError(f"Unsupported template file type: .{file_extension}")

