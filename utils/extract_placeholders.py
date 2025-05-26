import re
from io import BytesIO
from docx import Document # Import here to avoid circular import if needed elsewhere
from pptx import Presentation # Import here as well

def extract_placeholders(path: str) -> list[str]:
    """
    Return unique placeholder names (no braces) from a DOCX or PPTX file.
    Order is preserved from first appearance.
    """
    text_content = "" # Renamed for clarity

    if path.lower().endswith(".docx"):
        doc = Document(path)
        text_content = "\n".join(p.text for p in doc.paragraphs)
    elif path.lower().endswith(".pptx"): # Use elif for explicit handling
        prs = Presentation(path)
        texts_from_pptx = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame: # Check for text_frame to ensure it's a text-holding shape
                    texts_from_pptx.append(shape.text_frame.text)
        text_content = "\n".join(texts_from_pptx)
    # else: If an unsupported file type, text_content remains empty,
    # leading to an empty list of placeholders, which is acceptable behavior.

    # Extract {PLACEHOLDER} entries
    raw = re.findall(r"\{([^}]+)\}", text_content)
    seen = set()
    ordered = []
    for ph in raw:
        key = ph.strip()
        if key not in seen:
            seen.add(key)
            ordered.append(key)
    return ordered

def extract_text_from_file(stream: BytesIO, filename: str) -> str:
    """
    Extract text from uploaded file stream based on extension:
    - .txt  : decode UTF-8
    - .docx : extract all paragraphs
    - .pptx : extract all shape text from slides
    """
    ext = filename.lower().rsplit(".", 1)[-1]
    stream.seek(0) # Ensure stream is at the beginning

    if ext == "txt":
        return stream.read().decode("utf-8", errors="ignore")

    if ext == "docx":
        doc = Document(stream)
        return "\n".join(p.text for p in doc.paragraphs)

    if ext == "pptx":
        prs = Presentation(stream)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame: # Check for text_frame
                    texts.append(shape.text_frame.text)
        return "\n".join(texts)

    # fallback for unrecognized file types (e.g., if an image is passed here)
    # Attempt to decode as text, though it will likely be garbled for non-text files.
    return stream.read().decode("utf-8", errors="ignore")

