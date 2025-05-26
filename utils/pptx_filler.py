from pptx import Presentation
import pandas as pd, os

def generate_filled_pptx(template_path, outfile, uploads, text_inputs):
    prs = Presentation(template_path)

    # Map placeholder → replacement string (we convert everything to str in pptx)
    text_map = {}

    for ph, path in uploads.items():
        ext = os.path.splitext(path)[1].lower()
        if ext == ".txt":
            text_map[ph] = open(path).read()
        elif ext == ".docx":
            from docx import Document
            text_map[ph] = "\n".join(p.text for p in Document(path).paragraphs)
        elif ext in {".xlsx", ".csv"}:
            df = pd.read_excel(path) if ext == ".xlsx" else pd.read_csv(path)
            text_map[ph] = df.to_markdown(index=False)
        else:
            text_map[ph] = f"[{os.path.basename(path)}]"

    text_map.update(text_inputs)

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                for ph, rep in text_map.items():
                    if ph in shape.text:
                        shape.text = shape.text.replace(ph, str(rep))

    prs.save(outfile)
